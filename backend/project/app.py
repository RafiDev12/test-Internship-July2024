import os
import openai
import json
from pptx import Presentation
from pptx.util import Pt
from flask import Flask, request, render_template, send_file, jsonify
from flask_sqlalchemy import SQLAlchemy
from datetime import datetime

app = Flask(__name__)
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///presentations.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

openai.api_key = "YOUR_API_KEY"

class PresentationHistory(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(200), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    download_link = db.Column(db.String(200), nullable=False)

with app.app_context():
    db.create_all()

@app.route('/')
def index():
    history = PresentationHistory.query.order_by(PresentationHistory.created_at.desc()).all()
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        return render_template('history_table.html', history=history)
    return render_template('index.html', history=history)

@app.route('/generate', methods=['POST'])
def generate():
    request_data = request.get_json()
    presentation_title = request_data.get('title')

    if not presentation_title:
        return jsonify({"error": "Title is required"}), 400

    query_json = json.dumps({
        "input text": "[[QUERY]]",
        "output_format": "json",
        "json structure": {
            "slides": {
                "header": "string",
                "content": "string",
                "references": "string"
            }
        }
    }, indent=4)

    question = f"Generate a 10-slide presentation for the topic '{presentation_title}'. Each slide should have a header, detailed content, and references. Return as JSON."

    prompt = query_json.replace("[[QUERY]]", question)

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}]
    )

    response_content = response.choices[0].message['content']

    try:
        slide_data = json.loads(response_content)["slides"]
    except (json.JSONDecodeError, KeyError):
        return jsonify({"error": "Unable to parse the response from OpenAI"}), 500

    prs = Presentation()

    # Add the cover slide
    slide_layout = prs.slide_layouts[0]  # 0 index usually represents the cover slide layout
    cover_slide = prs.slides.add_slide(slide_layout)
    title = cover_slide.shapes.title
    title.text = presentation_title.upper()
    subtitle = cover_slide.placeholders[1]
    subtitle.text = "Generated by AI"

    # Add the content slides
    for slide in slide_data:
        slide_layout = prs.slide_layouts[1]  # 1 index represents a slide with title and content
        new_slide = prs.slides.add_slide(slide_layout)

        if 'header' in slide and slide['header']:
            title = new_slide.shapes.title
            title.text = slide['header']

        if 'content' in slide and slide['content']:
            shapes = new_slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide['content']
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(18)

        if 'references' in slide and slide['references']:
            # Add references as a new text box at the bottom of the slide
            left = Pt(50)
            top = Pt(450)
            width = Pt(860)
            height = Pt(100)
            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = "- " + slide['references']
            p.font.size = Pt(12)

    # Save the presentation in the PPTX folder
    if not os.path.exists('PPTX'):
        os.makedirs('PPTX')

    filename = os.path.join('PPTX', f"{presentation_title.replace(' ', '_')}_presentation.pptx")
    prs.save(filename)

    download_link = request.host_url + 'download/' + filename

    # Save to database
    new_presentation = PresentationHistory(title=presentation_title, download_link=download_link)
    db.session.add(new_presentation)
    db.session.commit()

    # Return success response with download link
    return jsonify({
        "message": "Presentation created successfully",
        "download_link": download_link
    }), 200

@app.route('/download/<path:filename>', methods=['GET'])
def download_file(filename):
    return send_file(filename, as_attachment=True)

@app.route('/search', methods=['GET'])
def search():
    query = request.args.get('query')
    if query:
        search_results = PresentationHistory.query.filter(PresentationHistory.title.ilike(f'%{query}%')).order_by(PresentationHistory.created_at.desc()).all()
    else:
        search_results = PresentationHistory.query.order_by(PresentationHistory.created_at.desc()).all()

    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        return render_template('history_table.html', history=search_results)

    return render_template('index.html', history=search_results)

if __name__ == '__main__':
    app.run(debug=True)
